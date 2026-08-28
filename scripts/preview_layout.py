# -*- coding: utf-8 -*-
"""販売図面PPTXを1枚のPNGに素描きして、レイアウトを目で確認できるようにする。

usage:
  python scripts/preview_layout.py 図面.pptx -o preview.png [--dpi 120]
  python scripts/preview_layout.py 図面.pptx --crop 0.4,10.0,10.7,15.6   # 部分拡大(cm)

LibreOffice が入っていない環境でも動く簡易プレビュー。
写真は埋め込み画像をそのまま描き、図形は塗り／枠線を、文字は同梱の Noto Serif JP で描く。
※ 本番フォントは HGS明朝E。字幅は全角1emで共通なので折返し位置の確認には使えるが、
   字形そのものの再現ではない。
"""
import argparse, io, os, sys
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import textfit

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
SERIF = os.path.join(ROOT, "assets", "fonts", "NotoSerifJP[wght].ttf")
SANS = os.path.join(ROOT, "assets", "fonts", "NotoSansJP[wght].ttf")
CM = 360000.0
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

_fcache = {}


def font(px, weight=400, serif=True):
    key = (px, weight, serif)
    if key not in _fcache:
        f = ImageFont.truetype(SERIF if serif else SANS, max(px, 4))
        try:
            f.set_variation_by_axes([weight])
        except Exception:
            pass
        _fcache[key] = f
    return _fcache[key]


def solid_color(el):
    """spPr 相当の要素から塗り色 (r,g,b) を拾う。取れなければ None。"""
    if el is None:
        return None
    fill = el.find(A + "solidFill")
    if fill is None:
        return None
    srgb = fill.find(A + "srgbClr")
    if srgb is None:
        return None
    v = srgb.get("val")
    try:
        return tuple(int(v[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return None


def run_color(r):
    rPr = r._r.find(A + "rPr")
    c = solid_color(rPr)
    return c or (30, 35, 40)


em_width = textfit.em
wrap = textfit.wrap


def draw_shape(sh, d, img, px_cm, ox, oy):
    if sh.left is None:
        return
    x, y = sh.left / CM - ox, sh.top / CM - oy
    w, h = (sh.width or 0) / CM, (sh.height or 0) / CM
    X, Y, X2, Y2 = (round(v * px_cm) for v in (x, y, x + w, y + h))

    if sh.shape_type == 13 and hasattr(sh, "image"):      # PICTURE
        try:
            im = Image.open(io.BytesIO(sh.image.blob)).convert("RGBA")
            im = im.resize((max(X2 - X, 1), max(Y2 - Y, 1)), Image.LANCZOS)
            img.paste(im, (X, Y), im)
        except Exception:
            d.rectangle([X, Y, X2, Y2], fill=(210, 205, 198))
        return

    if sh.has_table:
        draw_table(sh, d, px_cm, x, y, w)
        return

    spPr = sh._element.find(".//" + A + "spPr") if sh._element is not None else None
    fill = solid_color(spPr)
    ln = spPr.find(A + "ln") if spPr is not None else None
    line = solid_color(ln)
    if fill:
        d.rectangle([X, Y, X2, Y2], fill=fill)
    if line:
        d.rectangle([X, Y, X2, Y2], outline=line, width=max(1, round(0.02 * px_cm)))

    if sh.has_text_frame and sh.text_frame.text.strip():
        draw_text(sh, d, px_cm, x, y, w, h)


def draw_table(sh, d, px_cm, x, y, w):
    """表は罫線と文字だけを描く（塗りはテーマ依存なので省く）。"""
    tbl = sh.table
    cols = [c.width / CM for c in tbl.columns]
    cy = y
    for row in tbl.rows:
        rh = row.height / CM
        cx = x
        for ci, cell in enumerate(row.cells):
            cw = cols[ci]
            d.rectangle([round(cx * px_cm), round(cy * px_cm),
                         round((cx + cw) * px_cm), round((cy + rh) * px_cm)],
                        outline=(190, 175, 135), width=1)
            txt = cell.text.strip()
            if txt:
                pts = [r.font.size.pt for p in cell.text_frame.paragraphs
                       for r in p.runs if r.font.size]
                pt = max(pts) if pts else 9.0
                f = font(round(pt * 0.0352778 * px_cm), 600)
                d.text((round((cx + 0.1) * px_cm), round((cy + (rh - pt * 0.0352778) / 2) * px_cm)),
                       txt[:24], font=f, fill=(30, 35, 40))
            cx += cw
        cy += rh


def draw_text(sh, d, px_cm, x, y, w, h):
    tf = sh.text_frame
    bp = tf._txBody.bodyPr
    ins = {}
    for k, dflt in (("lIns", 0.254), ("rIns", 0.254), ("tIns", 0.127), ("bIns", 0.127)):
        v = bp.get(k)
        ins[k] = (int(v) / 914400 * 2.54) if v is not None else dflt
    no_wrap = bp.get("wrap") == "none"
    anchor = bp.get("anchor") or "t"
    tw = w - ins["lIns"] - ins["rIns"]

    lines = []
    for p in tf.paragraphs:
        if not p.runs:
            lines.append((0.0, "", 12.0, (30, 35, 40), 1.0, p.alignment))
            continue
        pt = max([(r.font.size.pt if r.font.size else 12.0) for r in p.runs])
        marl = 0.0
        if p._pPr is not None and p._pPr.get("marL"):
            marl = int(p._pPr.get("marL")) / 914400 * 2.54
        col = run_color(p.runs[0])
        mult = 1.2
        if p._pPr is not None:
            ls = p._pPr.find(A + "lnSpc")
            if ls is not None and len(ls) and ls[0].get("val"):
                mult = int(ls[0].get("val")) / 100000.0
        cap = 10 ** 6 if no_wrap else (tw - marl) / (pt * 0.0352778)
        for seg in wrap(p.text, cap):
            lines.append((marl, seg, pt, col, mult, p.alignment))

    total = sum(pt * 0.0352778 * mult * 1.2 for _, _, pt, _, mult, _ in lines)
    cy = y + ins["tIns"]
    if anchor == "ctr":
        cy = y + (h - total) / 2
    elif anchor == "b":
        cy = y + h - ins["bIns"] - total

    for marl, seg, pt, col, mult, algn in lines:
        lh = pt * 0.0352778 * mult * 1.2
        if seg:
            f = font(round(pt * 0.0352778 * px_cm), 600)
            tx = x + ins["lIns"] + marl
            if algn is not None and algn == 2:      # CENTER
                tx = x + (w - em_width(seg) * pt * 0.0352778) / 2
            d.text((round(tx * px_cm), round((cy + (lh - pt * 0.0352778) / 2) * px_cm)),
                   seg, font=f, fill=col)
        cy += lh


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out", default="preview.png")
    ap.add_argument("--dpi", type=float, default=120, help="1インチあたりのpx")
    ap.add_argument("--crop", help="切り出す範囲 left,top,right,bottom (cm)")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    px_cm = args.dpi / 2.54
    sw, sh_ = prs.slide_width / CM, prs.slide_height / CM
    ox = oy = 0.0
    if args.crop:
        l, t, r, b = (float(v) for v in args.crop.split(","))
        ox, oy, sw, sh_ = l, t, r - l, b - t
    img = Image.new("RGB", (round(sw * px_cm), round(sh_ * px_cm)), (255, 255, 255))
    d = ImageDraw.Draw(img)
    for shp in prs.slides[0].shapes:
        draw_shape(shp, d, img, px_cm, ox, oy)
    img.save(args.out)
    print(f"✅ {args.out}  ({img.width}x{img.height}px / {sw:.2f}x{sh_:.2f}cm)")


if __name__ == "__main__":
    main()
