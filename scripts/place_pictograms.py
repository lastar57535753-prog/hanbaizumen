# -*- coding: utf-8 -*-
"""販売図面PPTXの行頭にピクトグラムを差し込む。

usage:
  python scripts/place_pictograms.py 図面.pptx -o 図面_icon.pptx
  python scripts/place_pictograms.py 図面.pptx -o out.pptx --zones point,note,life
  python scripts/place_pictograms.py 図面.pptx -o out.pptx --color sumi --dry-run

考え方
------
「新しい意匠を勝手に足さない」（layout-rules.md）ので、**帯や枠は増やさない**。
POINT・備考は行頭の「・」をピクトグラムに**置き換える**だけ、
LIFE INFORMATION は施設名の左に1つ入れる。

  point : POINT本文の各行頭（10.5pt / 1行1枚）
  note  : 備考本文の各行頭（8.5pt）
  life  : LIFE INFORMATION の施設名の左（8.5pt）

字数が足りなくなる行が出たら**止める**（黙って溢れさせない）。
--force で警告のみにして続行できる。
"""
import argparse, os, sys
from pptx import Presentation
from pptx.util import Emu, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import textfit
from pick_pictograms import load_catalog, pick_line, LIFE_KEYS

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
CM = 360000.0
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
BULLETS = "・･●■▶◆〇○-–—"
MIN_ICON = 0.26          # これ以下は印刷で潰れるので使わない(cm)
FLOOR_PT = 8.0           # A4原寸で読める下限。これ以上は縮めない

# 見出しの文字列 → (ゾーン名, 走査するカタログ分類)
ZONE_HEADS = {
    "point": ("POINT", {"common", "unit", "site", "deal", "access"}),
    "note": ("備考", {"common", "unit", "site", "deal", "access"}),
    "life": ("LIFE INFORMATION", {"life", "access"}),
}


def cm_of(v):
    return (v or 0) / CM


def insets(tf):
    bp = tf._txBody.bodyPr
    out = {}
    for k, dflt in (("lIns", 0.254), ("rIns", 0.254), ("tIns", 0.127), ("bIns", 0.127)):
        v = bp.get(k)
        out[k] = (int(v) / 914400 * 2.54) if v is not None else dflt
    return out


def has_sp_autofit(tf):
    return tf._txBody.bodyPr.find(A + "spAutoFit") is not None


def para_pt(p, dflt=10.0):
    sizes = [r.font.size.pt for r in p.runs if r.font.size]
    return max(sizes) if sizes else dflt


def para_linespace(p):
    """lnSpc の倍率。指定が無ければ 1.0。"""
    if p._pPr is None:
        return 1.0
    ls = p._pPr.find(A + "lnSpc")
    if ls is None or not len(ls):
        return 1.0
    val = ls[0].get("val")
    return int(val) / 100000.0 if val else 1.0


def find_body(slide, head_text):
    """見出し図形の真下にある本文ボックスを返す。(見出し, 本文) を返す。

    同じ見出し語が複数ある図面（写真ギャラリー上の「POINT」など）もあるので、
    本文が見つかった見出しを採る。"""
    heads = [sh for sh in slide.shapes
             if sh.has_text_frame and sh.text_frame.text.strip() == head_text]
    for head in heads:
        hl = cm_of(head.left)
        hb = cm_of(head.top) + cm_of(head.height)
        hr = hl + cm_of(head.width)
        best = None
        for sh in slide.shapes:
            if sh is head or not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            l, t = cm_of(sh.left), cm_of(sh.top)
            if t < hb - 0.05 or t > hb + 1.2:          # 見出しのすぐ下
                continue
            if l < hl - 0.3 or l > hr:                  # 横位置が重なる
                continue
            if best is None or t < cm_of(best.top):
                best = sh
        if best is not None:
            return head, best
    return (heads[0] if heads else None), None


def room_below(slide, shape):
    """本文ボックスの下にどれだけ伸ばせるか(cm)。真下の図形までの距離。

    背景パネルや台紙のように本文ボックスを丸ごと含む図形は数えない。"""
    l, t = cm_of(shape.left), cm_of(shape.top)
    r, b = l + cm_of(shape.width), t + cm_of(shape.height)
    limit = None
    for sh in slide.shapes:
        if sh is shape or sh.left is None:
            continue
        sl, st = cm_of(sh.left), cm_of(sh.top)
        sr, sb = sl + cm_of(sh.width), st + cm_of(sh.height)
        if sl <= l + 0.05 and sr >= r - 0.05 and st <= t + 0.05 and sb >= b - 0.05:
            continue                                     # 台紙・背景
        if sr <= l + 0.05 or sl >= r - 0.05:
            continue                                     # 横に重ならない
        if st < b - 0.05:
            continue                                     # 下にない
        limit = st if limit is None else min(limit, st)
    return (limit - b) if limit is not None else 99.0


def line_layout(shape):
    """本文ボックスの各段落について (段落, 行数, 行の高さcm, 段落先頭のy) を返す。"""
    tf = shape.text_frame
    ins = insets(tf)
    w = cm_of(shape.width) - ins["lIns"] - ins["rIns"]
    no_wrap = tf._txBody.bodyPr.get("wrap") == "none"
    rows = []
    for p in tf.paragraphs:
        pt = para_pt(p)
        n = 1 if (no_wrap or not p.text.strip()) else textfit.lines_needed(p.text, w, pt)
        rows.append([p, n, pt, para_linespace(p)])
    total_lines = sum(r[1] for r in rows) or 1
    if has_sp_autofit(tf):
        # PowerPoint が実測して決めた高さがあるので、それを行数で割るのが一番確か
        lh_of = lambda pt, mult: (cm_of(shape.height) - ins["tIns"] - ins["bIns"]) / total_lines
    else:
        lh_of = lambda pt, mult: pt * mult * 1.2 * textfit.PT_CM
    y = cm_of(shape.top) + ins["tIns"]
    out = []
    for p, n, pt, mult in rows:
        lh = lh_of(pt, mult)
        out.append((p, n, lh, y, pt))
        y += lh * n
    return out, ins, w


def strip_bullet(p):
    for r in p.runs:
        if r.text and r.text[0] in BULLETS:
            r.text = r.text[1:].lstrip("　 ")
            return True
        if r.text.strip():
            return False
    return False


def set_indent(p, marl_cm):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(round(marl_cm * 914400 / 2.54))))
    pPr.set("indent", "0")


def widen_life_names(slide, name_box, log):
    """LIFE INFORMATION の距離欄を実測して詰め、施設名の欄をその分だけ広げる。

    施設名は枠いっぱいまで使っていることが多く、そのままではアイコンが入らない。
    距離欄（「徒歩 10分 （約800m）」）は右側に余りがあるので、そこを回す。"""
    nl, nt = cm_of(name_box.left), cm_of(name_box.top)
    nr = nl + cm_of(name_box.width)
    dist = None
    for sh in slide.shapes:
        if sh is name_box or not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        if abs(cm_of(sh.top) - nt) < 0.1 and cm_of(sh.left) >= nr - 0.1:
            if dist is None or cm_of(sh.left) < cm_of(dist.left):
                dist = sh
    if dist is None:
        return 0.0, None
    ins = insets(dist.text_frame)
    need = max(textfit.cm(p.text, para_pt(p)) for p in dist.text_frame.paragraphs
               if p.text.strip())
    free = cm_of(dist.width) - ins["lIns"] - ins["rIns"] - need - 0.05
    if free <= 0.05:
        return 0.0, dist
    dist.left = Emu(int((cm_of(dist.left) + free) * CM))
    dist.width = Emu(int((cm_of(dist.width) - free) * CM))
    name_box.width = Emu(int((cm_of(name_box.width) + free) * CM))
    log.append(f"   ※ 距離欄を {free:.2f}cm 詰めて施設名の欄を "
               f"{cm_of(name_box.width):.2f}cm に広げました")
    return free, dist


def set_pt(shape, pt):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                r.font.size = Pt(pt)


def shrink_to_make_room(body, log, mirror=None):
    """行数を増やさずにアイコン分の幅を空けるため、必要なら文字を小さくする。

    下限は FLOOR_PT（A4原寸で読める大きさ）。それでも空かなければ何もしない。
    ユーザーの決めごと「8pt を下回るなら縮小しない」に従う。"""
    tf = body.text_frame
    ins = insets(tf)
    text_w = cm_of(body.width) - ins["lIns"] - ins["rIns"]
    cur = max((para_pt(p) for p in tf.paragraphs if p.text.strip()), default=0)
    if not cur:
        return cur
    longest = max(textfit.em(p.text) for p in tf.paragraphs if p.text.strip())
    need = MIN_ICON * 1.18
    if text_w - longest * cur * textfit.PT_CM >= need:
        return cur                                   # そのままで入る
    pt = cur
    while pt - 0.5 >= FLOOR_PT:
        pt -= 0.5
        if text_w - longest * pt * textfit.PT_CM >= need:
            break
    else:
        return cur                                   # 8pt まで縮めても入らない
    set_pt(body, pt)
    if mirror is not None:
        set_pt(mirror, pt)          # 同じ行の距離欄も揃える
    log.append(f"   ※ 幅を作るため文字を {cur}→{pt}pt にしました（下限 {FLOOR_PT}pt・距離欄も同じ）")
    return pt


PICTO_PREFIX = "PICTO:"


def clear_pictograms(slide):
    """前回このスクリプトが置いたピクトグラムを消す（何度でも掛け直せるように）。"""
    n = 0
    for sh in list(slide.shapes):
        if (sh.name or "").startswith(PICTO_PREFIX):
            sh._element.getparent().remove(sh._element)
            n += 1
    return n


def place(slide, zone, items, color, force, dry, used, log):
    head_text, cats = ZONE_HEADS[zone]
    head, body = find_body(slide, head_text)
    if body is None:
        log.append(f"⚠ {zone}: 「{head_text}」の本文ボックスが見つかりません（スキップ）")
        return 0
    if zone == "life":
        _, dist_box = widen_life_names(slide, body, log)
        shrink_to_make_room(body, log, mirror=dist_box)
    rows, ins, text_w = line_layout(body)
    lh_min = min(r[2] for r in rows)
    body_pt = max(r[4] for r in rows)
    ideal = min(body_pt * textfit.PT_CM * 1.25, lh_min * 0.88)

    def extra_lines(reserve):
        """アイコン分だけ本文を狭めたときに増える行数と、溢れた行。"""
        w = text_w - reserve
        add, bad = 0, []
        for p, n, lh, y, pt in rows:
            t = p.text.lstrip(BULLETS).lstrip("　 ")
            if not t:
                continue
            got = textfit.lines_needed(t, w, pt)
            if got > n:
                add += got - n
                bad.append((t, textfit.em(t), textfit.capacity(w, pt)))
        return add, bad

    # 行が増えても、本文ボックスの下に伸びしろがあるなら許す。
    # 伸ばせない場合はアイコンを縮めて収める。それも無理なら止める。
    room = 0.0 if zone == "life" else room_below(slide, body)
    size, gap = ideal, ideal * 0.18
    while True:
        reserve = size + gap
        add, bad = extra_lines(reserve)
        if add == 0 or add * lh_min <= room:
            break
        size -= 0.02
        gap = size * 0.18
        if size < MIN_ICON:
            for t, e, cap in bad:
                log.append(f"   ✗ 行が増える: {e:.1f}字 > {cap:.1f}字  「{t[:30]}」")
            log.append(f"   （下の余白は {room:.2f}cm しかなく、{MIN_ICON}cm 未満のアイコンは小さすぎます）")
            if not force:
                raise SystemExit(
                    f"!! {zone}: アイコンを入れる幅がありません（本文が枠の端まで使っています）。\n"
                    f"   本文を1〜2字短くするか、--force で承知のうえ続行してください。")
            size, gap = ideal, ideal * 0.18
            break
    reserve = size + gap
    add, _ = extra_lines(reserve)
    if add:
        log.append(f"   ※ 本文が {add} 行増えます（下に {room:.2f}cm の余白あり）")
    if size < ideal - 0.001:
        log.append(f"   ※ 幅が足りないのでアイコンを {ideal:.2f}→{size:.2f}cm に縮めました")

    n_placed = 0
    for p, n, lh, y, pt in rows:
        text = p.text
        if not text.strip():
            continue
        it, found = pick_line(text, items, cats=cats, used=used)
        if it is None:
            log.append(f"   ・該当なし: 「{text[:26]}」")
            continue
        png = os.path.join(ROOT, "pictograms", "png", color, it["key"] + ".png")
        if not os.path.exists(png):
            log.append(f"   ✗ 画像なし: {png}")
            continue
        used.add(it["key"])
        x = cm_of(body.left) + ins["lIns"]
        iy = y + (lh - size) / 2
        log.append(f"   ✅ {it['label']:<10} {found[0]:<10} ({x:.2f}, {iy:.2f}) {size:.2f}cm"
                   f"  ← 「{text[:22]}」")
        if not dry:
            pic = slide.shapes.add_picture(png, Emu(int(x * CM)), Emu(int(iy * CM)),
                                           Emu(int(size * CM)), Emu(int(size * CM)))
            pic.name = PICTO_PREFIX + it["key"]
            strip_bullet(p)
            set_indent(p, ins["lIns"] + reserve)
        n_placed += 1
    return n_placed


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out")
    ap.add_argument("--zones", default="point,note",
                    help="point / note / life をカンマ区切りで（既定: point,note）")
    ap.add_argument("--color", default="gold", choices=["sumi", "gold", "gold_light", "white"])
    ap.add_argument("--dry-run", action="store_true", help="書き込まずに配置予定だけ出す")
    ap.add_argument("--force", action="store_true", help="行が増える場合も続行する")
    args = ap.parse_args()
    if not args.dry_run and not args.out:
        ap.error("-o を指定してください（--dry-run なら不要）")

    prs = Presentation(args.pptx)
    slide = prs.slides[0]
    items = load_catalog()["items"]
    if not args.dry_run:
        gone = clear_pictograms(slide)
        if gone:
            print(f"■ 前回のピクトグラム {gone}枚を差し替えます")
    used, total = set(), 0
    for zone in [z.strip() for z in args.zones.split(",") if z.strip()]:
        if zone not in ZONE_HEADS:
            sys.exit(f"!! 未知のゾーン: {zone}")
        log = []
        print(f"■ {zone}（{ZONE_HEADS[zone][0]}）")
        try:
            n = place(slide, zone, items, args.color, args.force, args.dry_run, used, log)
        finally:
            for l in log:
                print(l)
        print(f"   → {n}枚")
        total += n

    if args.dry_run:
        print(f"\n（--dry-run のため書き込んでいません。合計 {total}枚）")
        return
    prs.save(args.out)
    print(f"\n✅ {total}枚を配置 → {args.out}")


if __name__ == "__main__":
    main()
