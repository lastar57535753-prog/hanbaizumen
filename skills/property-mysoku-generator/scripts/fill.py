# -*- coding: utf-8 -*-
"""穴あきテンプレに値を流し込んで販売図面pptxを作る。

usage: python fill.py <template.pptx> <data.json> <out.pptx> [--agent キー]

data.json のキーはプレースホルダ名（{{}} なし）。
値が無いキーは【要確認】を赤字で残す。

担当者（AGENT_NAME / AGENT_MOBILE / AGENT_MAIL）は assets/agents.json から補う。
data.json に "AGENT": "iwasawa" と書くか --agent iwasawa で切り替わる。
data.json に直接 AGENT_* を書けばそちらが優先。
会社名・住所・代表TEL・免許番号は共通なのでテンプレート側に固定してある。
"""
import sys, io, json, os, re
from pptx import Presentation
from pptx.dml.color import RGBColor

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

TODO = "【要確認】"
RED = RGBColor(0xFF, 0x00, 0x00)
PH = re.compile(r"\{\{([A-Z0-9_]+)\}\}")

# 「無いのが普通」の枠。値が無ければ【要確認】ではなく空欄にする。
# BUILDER2 = 施工がJVのときの2社目。TERRACE = テラス/専用庭付き住戸のみ。
OPTIONAL = {"BUILDER2", "TERRACE", "BADGE1", "BADGE2", "SUBCATCH", "DIRECTION", "AREA_MINI",
            "CATCH2",
            "POINT4", "POINT5", "POINT6", "POINT7", "POINT8",
            "NOTE1", "NOTE2", "NOTE3", "NOTE4",
            "CAP1", "CAP2", "CAP3", "CAP4", "CAP5", "CAP6"}


def apply_agent(data, explicit=None):
    """担当者の連絡先を data に補う。既に data にある値は上書きしない。"""
    path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                        "assets", "agents.json")
    if not os.path.exists(path):
        return data
    with open(path, encoding="utf-8") as f:
        reg = json.load(f)
    key = explicit or data.get("AGENT") or reg.get("default")
    who = reg.get("agents", {}).get(key)
    if who is None:
        sys.exit(f"!! 担当者 '{key}' が assets/agents.json にありません"
                 f"（登録済み: {'・'.join(reg.get('agents', {}))}）")
    for k, v in who.items():
        data.setdefault(k, v)
    print(f"担当: {who.get('AGENT_FULLNAME', key)}")
    return data


def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        try:
            if sh.shape_type == 6:
                yield from iter_shapes(sh.shapes)
        except Exception:
            pass


def fill_textframe(tf, data, acc):
    filled, todos, blanks = acc
    for p in tf.paragraphs:
        for r in p.runs:
            keys = PH.findall(r.text)
            if not keys:
                continue
            missing = False
            out = r.text
            for key in keys:
                val = data.get(key)
                if val is None or val == "":
                    if key in OPTIONAL:
                        out = out.replace("{{%s}}" % key, "")
                        blanks.append(key)
                    else:
                        out = out.replace("{{%s}}" % key, TODO)
                        todos.append(key)
                        missing = True
                else:
                    out = out.replace("{{%s}}" % key, str(val))
                    filled[0] += 1
            r.text = out
            if missing:
                r.font.color.rgb = RED


def fill(prs, data):
    acc = ([0], [], [])   # filled(mutable), todos, blanks
    for sh in iter_shapes(prs.slides[0].shapes):
        if sh.has_text_frame:
            fill_textframe(sh.text_frame, data, acc)
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    fill_textframe(cell.text_frame, data, acc)
    return acc[0][0], acc[1], acc[2]


def main():
    tpl, datafile, out = sys.argv[1], sys.argv[2], sys.argv[3]
    rest = sys.argv[4:]
    agent = rest[rest.index("--agent") + 1] if "--agent" in rest else None
    with open(datafile, encoding="utf-8") as f:
        data = json.load(f)
    data = apply_agent(data, agent)
    prs = Presentation(tpl)
    filled, todos, blanks = fill(prs, data)
    prs.save(out)
    print(f"filled: {filled}")
    if todos:
        print(f"【要確認】として赤字で残した項目 ({len(todos)}):")
        for k in todos:
            print(f"  - {k}")
    if blanks:
        print(f"任意項目のため空欄にした ({len(blanks)}): {', '.join(blanks)}")
    print(f"saved: {out}")


if __name__ == "__main__":
    main()
