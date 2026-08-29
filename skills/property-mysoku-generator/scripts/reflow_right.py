# -*- coding: utf-8 -*-
"""右カラム（物件概要→備考→LIFE INFORMATION）を、各ブロックの高さで縦に積み直す。
順序・左右位置・幅は固定。縦幅だけ物件ごとに可変にして重なりを防ぐ。

usage: python reflow_right.py <in.pptx> <out.pptx>

方式（自動選択）:
  - Windows(PowerPoint) → COM でテーブルの折り返し込み"実描画高さ"を使う（最も正確）
  - macOS / Linux       → python-pptx のみで stored/推定高さから積み直す（要目視確認）

テンプレ固定ID:
  305=概要ヘッダ 307=概要テーブル 311=備考ヘッダ 313=備考本文
  308=LIFEヘッダ 327=LIFE左列 328=LIFE右列   146=フッター帯(下限)
"""
import sys, io, os

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

PP_AUTOSIZE_FIT = 1   # ppAutoSizeShapeToFitText
GAP = 5.0             # ブロック間の余白(pt)


def _reflow_com(src, out):
    import win32com.client
    app = win32com.client.Dispatch("PowerPoint.Application")
    pres = app.Presentations.Open(src, WithWindow=False)
    try:
        slide = pres.Slides(1)
        S = {sh.Id: sh for sh in slide.Shapes}
        tbl, bikoh, biko, lifeh, lifeL, lifeR, footer = (
            S[307], S[311], S[313], S[308], S[327], S[328], S.get(146))
        qlabel, qbox = S.get(21), S.get(22)   # ↓物件詳細↓ ラベル と QR枠

        foot_top = footer.Top if footer else 535.0
        try:
            biko.TextFrame.AutoSize = PP_AUTOSIZE_FIT
        except Exception:
            pass

        table_bottom = tbl.Top + tbl.Height
        print(f"table: top={tbl.Top:.1f} height={tbl.Height:.1f} bottom={table_bottom:.1f} (footer_top={foot_top:.1f})")

        y = table_bottom + GAP
        bikoh.Top = y
        yb = bikoh.Top + bikoh.Height
        qbot = yb
        if qlabel is not None and qbox is not None:
            qlabel.Top = bikoh.Top
            qbox.Top = qlabel.Top + qlabel.Height
            qbot = qbox.Top + qbox.Height
            biko.Width = max(120.0, qlabel.Left - biko.Left - 4)
        biko.Top = yb
        ybody = biko.Top + biko.Height
        y = max(ybody, qbot) + GAP
        lifeh.Top = y;                 y = lifeh.Top + lifeh.Height
        lifeL.Top = y
        lifeR.Top = y;                 y = y + max(lifeL.Height, lifeR.Height)
        print(f"備考ヘッダ={bikoh.Top:.1f}  備考本文={biko.Top:.1f}(h={biko.Height:.1f})  QR下={qbot:.1f}  LIFEヘッダ={lifeh.Top:.1f}  末尾={y:.1f}")

        if y > foot_top:
            print(f"⚠ フッター({foot_top:.1f})を{y-foot_top:.1f}pt超過 — 要フォント調整")
        else:
            print(f"OK: フッターまで余裕 {foot_top-y:.1f}pt")

        pres.SaveAs(out)
        print("saved:", out)
    finally:
        pres.Close(); app.Quit()


def _reflow_pptx(src, out):
    """COM が無い環境(mac/Linux)向け。python-pptx の stored/推定高さで積み直す。
    テーブルの折り返しは正確には測れないので、備考本文は行数から高さを推定する。
    出力後は PowerPoint/LibreOffice で目視確認すること。"""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    PT = 12700  # 1pt = 12700 EMU
    prs = Presentation(src)
    slide = prs.slides[0]
    S = {sh.shape_id: sh for sh in slide.shapes}

    def top(s): return s.top
    def bottom(s): return s.top + s.height

    tbl = S.get(307); bikoh = S.get(311); biko = S.get(313)
    lifeh = S.get(308); lifeL = S.get(327); lifeR = S.get(328); footer = S.get(146)
    qlabel, qbox = S.get(21), S.get(22)
    if not all([tbl, bikoh, biko, lifeh, lifeL, lifeR]):
        # 必須ブロックが無ければ触らずコピー
        prs.save(out); print("reflow(pptx): 必須ID不足のため無調整で保存:", out); return

    foot_top = top(footer) if footer is not None else int(535.0 * PT)
    gap = int(GAP * PT)

    # 備考本文の高さを行数×行送りで推定（stored高さが実態と乖離しやすいため）
    def est_text_height(shape):
        tf = shape.text_frame
        total = 0
        for p in tf.paragraphs:
            txt = p.text
            # フォントサイズ(pt)を推定
            sz = None
            if p.runs and p.runs[0].font.size:
                sz = p.runs[0].font.size.pt
            sz = sz or 9.0
            n = max(1, txt.count("\n") + 1) if txt.strip() else 1
            total += int(n * sz * 1.35 * PT)
        # 上下パディング少々
        return total + int(6 * PT)

    try:
        biko_h = est_text_height(biko)
        biko.height = biko_h
    except Exception:
        biko_h = biko.height

    y = bottom(tbl) + gap
    bikoh.top = y
    yb = bikoh.top + bikoh.height
    qbot = yb
    if qlabel is not None and qbox is not None:
        qlabel.top = bikoh.top
        qbox.top = qlabel.top + qlabel.height
        qbot = qbox.top + qbox.height
        biko.width = max(int(120 * PT), qlabel.left - biko.left - int(4 * PT))
    biko.top = yb
    ybody = biko.top + biko.height
    y = max(ybody, qbot) + gap
    lifeh.top = y; y = lifeh.top + lifeh.height
    lifeL.top = y; lifeR.top = y
    y = y + max(lifeL.height, lifeR.height)

    over = y - foot_top
    if over > 0:
        print(f"⚠ フッターを約{over/PT:.1f}pt超過の見込み（推定）— フォント調整 or 目視確認")
    else:
        print(f"OK(推定): フッターまで約{-over/PT:.1f}pt")
    prs.save(out)
    print("saved:", out, "（reflow=python-pptx推定。PowerPoint/LibreOfficeで目視確認推奨）")


def main():
    src = os.path.abspath(sys.argv[1])
    out = os.path.abspath(sys.argv[2]) if len(sys.argv) > 2 else src
    try:
        import win32com.client  # noqa: F401
        has_com = True
    except Exception:
        has_com = False
    if has_com:
        _reflow_com(src, out)
    else:
        print("ℹ COM無し(mac/Linux)→ python-pptxで右カラムを推定整列します。")
        _reflow_pptx(src, out)


if __name__ == "__main__":
    main()
