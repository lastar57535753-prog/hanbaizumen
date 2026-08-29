# -*- coding: utf-8 -*-
"""出力前セルフQC：被り・はみ出し・整列・文字あふれ・書体違反を検査する安全網。

usage: python preflight.py <final.pptx> [--warn-only]
終了コード: 0=合格 / 1=要確認あり（--warn-only なら常に0）

計測は2通り。
  Windows + PowerPoint … COM の実描画座標（いちばん正確）
  それ以外            … python-pptx + 全角1em換算（pptx_geom）

**どちらの環境でも検査して、違反があれば終了コード1で止まる。**
以前は Windows 以外だと無検査で合格を返していたので、mac/Linux では
安全網がまったく効いていなかった。
"""
import io
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

PT = 28.3465  # 1cm あたりのポイント数（COM の座標単位）
SLIDE_W = 29.69

# 重なってはいけない主要コンテンツ（テンプレ固定ID）
CONTENT = {25: "キャッチ", 216: "POINT本文", 292: "POINT見出し", 287: "価格", 281: "アクセス",
           277: "物件名", 307: "物件概要表", 311: "備考見出し", 313: "備考本文",
           308: "LIFE見出し", 327: "LIFE左", 328: "LIFE右", 21: "物件詳細",
           85: "タグ:リビング", 34: "タグ", 35: "タグ", 267: "タグ", 268: "タグ",
           269: "タグ", 272: "タグ", 314: "連絡先", 318: "案内文"}
ALIGN_LEFT = [25, 292, 216]
FOOTER_IDS = {314, 318, 321, 322, 325, 326}
WHITELIST = {frozenset((292, 216)), frozenset((311, 21)), frozenset((311, 313)),
             frozenset((305, 307)), frozenset((308, 327)), frozenset((308, 328)),
             frozenset((314, 318))}
ALLOWED_FONT = "HGS明朝E"      # 図面で使ってよい書体はこれだけ
# 主要画像（間取図・写真）に重なってはいけないヘッダー文字
OVER_IMG = {25: "キャッチ", 287: "価格"}


def overlap_area(a, b):
    ix = max(0, min(a[2], b[2]) - max(a[0], b[0]))
    iy = max(0, min(a[3], b[3]) - max(a[1], b[1]))
    return ix * iy


def check(shapes, footer_top, badge_rect=None, pics=None):
    """shapes: {id: dict(name, box, text, height_cm, text_h_cm, autosize, fonts)}"""
    issues = []
    present = {i: s for i, s in shapes.items() if i in CONTENT}

    ids = list(present)
    for a in range(len(ids)):
        for b in range(a + 1, len(ids)):
            if frozenset((ids[a], ids[b])) in WHITELIST:
                continue
            ov = overlap_area(present[ids[a]]["text"], present[ids[b]]["text"])
            if ov > 0.15:
                issues.append(f"被り: 「{CONTENT[ids[a]]}」×「{CONTENT[ids[b]]}」 {ov:.2f}cm²")
    # キャッチ・価格が画像の**下敷き**になっていないか。
    # 重なり自体は違反ではない（キャッチはスクリム越しにヒーロー写真の上に載る意匠）。
    # 問題になるのは絵が文字より前面にあるとき＝文字が隠れるときだけなので z順で見る。
    for i, label in OVER_IMG.items():
        if i not in present:
            continue
        for pb, pz in (pics or []):
            if pz < present[i]["z"]:
                continue
            ov = overlap_area(present[i]["text"], pb)
            if ov > 0.15:
                issues.append(f"被り: 「{label}」が画像の下に隠れています {ov:.2f}cm²")
                break

    if badge_rect and 85 in present:
        ov = overlap_area(present[85]["box"], badge_rect)
        if ov > 0.15:
            issues.append(f"被り: 「タグ:リビング」×「高級バッジ」 {ov:.2f}cm²")

    for i, s in present.items():
        if i in FOOTER_IDS:
            continue
        r = s["text"]
        if r[3] > footer_top + 0.05:
            issues.append(f"はみ出し: 「{CONTENT[i]}」がフッター({footer_top:.1f}cm)を"
                          f"{r[3]-footer_top:.2f}cm超過")
        # 自動リサイズ図形（wrap=none / 自動調整）こそ枠が文字に追従して紙面外へ
        # 伸びるので、枠ではなく**文字の実寸**で左右を見る（価格が長いと右へ抜ける）
        if r[0] < -0.1 or r[2] > SLIDE_W + 0.11:
            issues.append(f"はみ出し: 「{CONTENT[i]}」が紙面左右外へ")

    lefts = [(i, present[i]["box"][0]) for i in ALIGN_LEFT if i in present]
    if len(lefts) >= 2:
        mn = min(l for _, l in lefts)
        mx = max(l for _, l in lefts)
        if mx - mn > 0.08:
            issues.append(f"整列: 左端が不揃い（{[CONTENT[i] for i,_ in lefts]} 差{(mx-mn)*10:.1f}mm）")

    for i, s in present.items():
        if s["autosize"] or not s["text_h"]:
            continue
        if s["text_h"] > s["h"] + 0.06:
            issues.append(f"文字あふれ(縦): 「{CONTENT[i]}」 文字{s['text_h']:.2f}>枠{s['h']:.2f}cm")

    bad = sorted({f for s in shapes.values() for f in s["fonts"] if f != ALLOWED_FONT})
    if bad:
        where = sorted({CONTENT.get(i, f"id{i}") for i, s in shapes.items()
                        if any(f != ALLOWED_FONT for f in s["fonts"])})
        issues.append(f"書体違反: {'・'.join(bad)} が使われています（{'・'.join(where)}）"
                      f"／全文字 {ALLOWED_FONT} にすること")
    return issues


def collect_com(path):
    import win32com.client
    app = win32com.client.Dispatch("PowerPoint.Application")
    pres = app.Presentations.Open(os.path.abspath(path), WithWindow=False)
    try:
        sl = pres.Slides(1)
        shapes, badge, footer_top, pics = {}, None, 21.0, []
        for sh in sl.Shapes:
            if sh.Type == 13 and 3 * PT <= sh.Width < 25 * PT:
                pics.append(((sh.Left / PT, sh.Top / PT,
                              (sh.Left + sh.Width) / PT, (sh.Top + sh.Height) / PT),
                             sh.ZOrderPosition))
            if sh.Name == "LUXBADGE_PLATE":
                badge = (sh.Left / PT, sh.Top / PT,
                         (sh.Left + sh.Width) / PT, (sh.Top + sh.Height) / PT)
            if abs(sh.Width / PT - SLIDE_W) < 1 and sh.Top / PT > 17.5:
                footer_top = min(footer_top, sh.Top / PT)
            b = (sh.Left / PT, sh.Top / PT,
                 (sh.Left + sh.Width) / PT, (sh.Top + sh.Height) / PT)
            t, th, fonts = b, 0.0, set()
            try:
                if sh.HasTextFrame and sh.TextFrame.HasText:
                    tr = sh.TextFrame.TextRange
                    t = (tr.BoundLeft / PT, tr.BoundTop / PT,
                         (tr.BoundLeft + tr.BoundWidth) / PT, (tr.BoundTop + tr.BoundHeight) / PT)
                    th = tr.BoundHeight / PT
                    fonts.add(tr.Font.NameFarEast or tr.Font.Name)
            except Exception:
                pass
            shapes[sh.Id] = {"box": b, "text": t, "h": sh.Height / PT, "text_h": th,
                             "autosize": False, "fonts": {f for f in fonts if f},
                             "z": sh.ZOrderPosition}
        return shapes, footer_top, badge, pics
    finally:
        pres.Close()
        app.Quit()


def collect_pptx(path):
    import pptx_geom as g
    from pptx import Presentation
    sl = Presentation(path).slides[0]
    shapes, badge, footer_top, pics = {}, None, 21.0, []
    for z, sh in enumerate(sl.shapes):
        if sh.left is None:
            continue
        b = g.box(sh)
        if sh.shape_type == 13 and 3 <= g.cm_of(sh.width) < 25:
            pics.append((b, z))
        if (sh.name or "") == "LUXBADGE_PLATE":
            badge = b
        if abs(g.cm_of(sh.width) - SLIDE_W) < 1 and b[1] > 17.5:
            footer_top = min(footer_top, b[1])
        auto = False
        if sh.has_text_frame:
            bp = sh.text_frame._txBody.bodyPr
            auto = (bp.get("wrap") == "none"
                    or bp.find(g.A + "spAutoFit") is not None
                    or bp.find(g.A + "normAutofit") is not None)
        shapes[sh.shape_id] = {"box": b, "text": g.text_rect(sh), "h": g.cm_of(sh.height),
                               "text_h": g.text_height(sh), "autosize": auto,
                               "fonts": g.fonts(sh), "z": z}
    return shapes, footer_top, badge, pics


def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    warn_only = "--warn-only" in sys.argv
    if not args:
        sys.exit("usage: python preflight.py <final.pptx> [--warn-only]")
    path = args[0]
    try:
        import win32com.client  # noqa: F401
        shapes, footer_top, badge, pics = collect_com(path)
        how = "PowerPoint(COM)の実描画座標"
    except Exception:
        shapes, footer_top, badge, pics = collect_pptx(path)
        how = "python-pptx＋全角1em換算（PowerPointが無い環境）"

    issues = check(shapes, footer_top, badge, pics)
    print(f"── 出力前セルフQC（{how}）──")
    if not issues:
        print("✅ 合格：被り・はみ出し・整列・文字あふれ・書体 いずれも問題なし（QRのみ手動）")
    else:
        print(f"⚠ 要確認 {len(issues)}件：")
        for m in issues:
            print("   -", m)
        if "COM" not in how:
            print("   ※ この環境の計測は近似です。最終確認は PowerPoint で開いて行ってください。")
    sys.exit(0 if (warn_only or not issues) else 1)


if __name__ == "__main__":
    main()
