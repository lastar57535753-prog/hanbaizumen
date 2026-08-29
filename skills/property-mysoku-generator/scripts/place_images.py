# -*- coding: utf-8 -*-
"""穴あきテンプレに実画像を挿入し、レイアウトを整える。

usage: python place_images.py <in.pptx> <config.json> <out.pptx>

config.json:
{
  "images": {
     "116": {"path": "...jpg", "mode": "fill"},   # fill=中央クロップで枠を埋める（写真）
     "139": {"path": "...png", "mode": "fit"}      # fit=縦横比維持で枠内に収める（間取図）
  },
  "relabels": { "307": { "5,0": "バルコニー面積" } },  # テーブルのラベル差し替え(行,列)
  "point_fix": {"id": 216, "bottom_cm": 15.0}         # POINTが写真にかぶらないよう箱を固定
}

- 画像は元の図形と同じ位置・zオーダーに差し込む（キャプションを隠さない）。
- fill: 枠比率に合わせて中央クロップ（歪ませない）。fit: 収める（間取図＝引き伸ばし厳禁）。
"""
import sys, io, json, os, re, copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image, ImageOps, ImageFilter, ImageEnhance

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import textfit
import hashlib
import urllib.request

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
CM = 360000




# ── 写真の取得元 ─────────────────────────────────────────────────
# images.json の "path" にはローカルパスのほか、Dropbox等の共有URLも書ける。
# 共有URLは1度だけ落としてキャッシュし、2回目からはローカルを使う。
_DL_CACHE = None


def _cache_dir(out_path):
    global _DL_CACHE
    if _DL_CACHE is None:
        _DL_CACHE = os.path.join(os.path.dirname(os.path.abspath(out_path)), "_dl")
        os.makedirs(_DL_CACHE, exist_ok=True)
    return _DL_CACHE


def _direct_url(url):
    """共有リンクを「画像そのもの」が返るURLに直す。"""
    if "dropbox.com" in url:
        u = re.sub(r"[?&]dl=\d", "", url)
        u = re.sub(r"[?&]raw=\d", "", u)
        return u + ("&" if "?" in u else "?") + "dl=1"
    if "drive.google.com" in url:
        m = re.search(r"/d/([\w-]+)", url) or re.search(r"[?&]id=([\w-]+)", url)
        if m:
            return f"https://drive.google.com/uc?export=download&id={m.group(1)}"
    return url


def resolve_path(path, out_path, allow_missing=False):
    """images.json の path を実ファイルに解決する。URLなら取得してキャッシュする。

    取れなかったら**止める**。写真が抜けたまま図面が出来てしまうのがいちばん困るため
    （以前は警告だけ出して先へ進み、写真ゼロの図面が納品されかけた）。
    どうしても穴を許したいときだけ --allow-missing。"""
    if not str(path).lower().startswith(("http://", "https://")):
        if os.path.exists(path):
            return path
        if allow_missing:
            return None
        sys.exit(f"!! 画像が見つかりません: {path}\n"
                 f"   パスを直すか、承知のうえで進めるなら --allow-missing を付けてください。")

    url = _direct_url(path)
    name = hashlib.sha1(url.encode("utf-8")).hexdigest()[:16]
    for ext in (".jpg", ".jpeg", ".png", ".webp"):
        cached = os.path.join(_cache_dir(out_path), name + ext)
        if os.path.exists(cached):
            return cached
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=60) as r:
            blob = r.read()
            ctype = (r.headers.get("Content-Type") or "").lower()
    except Exception as e:
        if allow_missing:
            print(f"⚠ 取得できません: {path} ({e})")
            return None
        sys.exit(f"!! 画像を取得できません: {path}\n   {e}\n"
                 f"   共有リンクが「閲覧可能」になっているか、ネットワークの許可設定を確認してください。")
    ext = ".png" if "png" in ctype else ".webp" if "webp" in ctype else ".jpg"
    dst = os.path.join(_cache_dir(out_path), name + ext)
    with open(dst, "wb") as f:
        f.write(blob)
    try:
        Image.open(dst).verify()
    except Exception:
        os.remove(dst)
        if allow_missing:
            print(f"⚠ 画像として読めません: {path}")
            return None
        sys.exit(f"!! 取得したデータが画像ではありません: {path}\n"
                 f"   Dropboxならファイルへの共有リンク（フォルダではなく）を使ってください。")
    print(f"取得: {os.path.basename(dst)} <- {path[:60]}")
    return dst


# 被写体の種類ごとのクロップ・アンカー（縦位置 0=上 1=下）
ANCHOR = {"exterior": 0.12, "view": 0.45, "room": 0.5, "living": 0.5,
          "kitchen": 0.42, "bath": 0.5, "hall": 0.5, "facility": 0.5}
# 自動ハイブリッドのしきい値（クロップ損失%）
CROP_OK = 20     # これ以下 → 中央クロップ
ANCHOR_OK = 40   # これ以下 → 賢いアンカークロップ / 超えたら → ぼかし背景フィル


def print_px(cm, dpi=300):
    return max(1, round(cm / 2.54 * dpi))


def blurred_fill(img, W, H):
    """全体を収めて、余白は同じ写真の拡大ぼかしで埋める（切れ0・歪み0・白帯0）。"""
    bg = ImageOps.fit(img, (W, H), Image.LANCZOS)
    bg = bg.filter(ImageFilter.GaussianBlur(max(6, W // 55)))
    bg = ImageEnhance.Brightness(bg).enhance(0.82)
    fg = ImageOps.contain(img, (W, H), Image.LANCZOS)
    out = bg.copy()
    out.paste(fg, ((W - fg.width) // 2, (H - fg.height) // 2))
    return out


def smart_composite(path, wcm, hcm, kind, tmpdir):
    """枠比率に合わせて自動ハイブリッド合成。(合成ファイル, 方式, 損失%, 実効dpi) を返す。"""
    img = Image.open(path).convert("RGB")
    Wp, Hp = print_px(wcm), print_px(hcm)
    Ri, Rf = img.width / img.height, Wp / Hp
    loss = (1 - min(Ri, Rf) / max(Ri, Rf)) * 100
    if loss <= CROP_OK:
        method, cy = "中央クロップ", 0.5
    elif loss <= ANCHOR_OK:
        method, cy = "アンカークロップ", ANCHOR.get(kind, 0.5)
    else:
        method, cy = "ぼかし背景フィル", None
    out = blurred_fill(img, Wp, Hp) if cy is None else ImageOps.fit(
        img, (Wp, Hp), Image.LANCZOS, centering=(0.5, cy))
    eff_dpi = round(min(img.width / (wcm / 2.54), img.height / (hcm / 2.54)))
    os.makedirs(tmpdir, exist_ok=True)
    fn = os.path.join(tmpdir, f"cmp_{os.path.basename(path)}.jpg")
    out.save(fn, quality=92)
    return fn, method, round(loss), eff_dpi


JP_OPEN = "「（【(『［｛《〈"
JP_CLOSE = "」）】)』］｝》〉"
NO_START = "、。，．・）」』】）,.:：;；!！?？ー〜%"     # 行頭に置けない（前行へ）
_ASCII = re.compile(r"[0-9A-Za-z%．.,／/－ー~〜:：＋+#&']+")
_KATA = re.compile(r"[ァ-ヶー・]{2,}")   # カタカナ語（ワイドサッシ/コンシェルジュ等）は分割しない


def _tokenize(t):
    """語中で切らないための最小トークン化。英数字連続・カタカナ語・『「…」』等の括弧内は分割不可。"""
    toks, i, n = [], 0, len(t)
    while i < n:
        c = t[i]
        if c in JP_OPEN:
            close = JP_CLOSE[JP_OPEN.index(c)]
            j = t.find(close, i + 1)
            j = j if j != -1 else i
            toks.append(t[i:j + 1]); i = j + 1; continue
        m = _ASCII.match(t, i)
        if m and m.end() > i:
            toks.append(t[i:m.end()]); i = m.end(); continue
        m = _KATA.match(t, i)
        if m and m.end() > i:
            toks.append(t[i:m.end()]); i = m.end(); continue
        toks.append(c); i += 1
    return toks


def jp_wrap(text, max_cm, pt):
    """日本語を語中で切らずに折り返す。禁則（行頭約物）も簡易対応。

    幅は「全角1em・半角0.5em」で数える（textfit）。実フォントで計測しないのは、
    本番の HGS明朝E が mac/Linux に無く、環境ごとに折返し位置が変わってしまうため。
    和文は全角1emなので、この数え方なら Windows でも mac でも同じ結果になる。"""
    lines, cur = [], ""
    for tok in _tokenize(text):
        cand = cur + tok
        if cur and textfit.cm(cand, pt) > max_cm:
            lines.append(cur); cur = tok
        else:
            cur = cand
    if cur:
        lines.append(cur)
    fixed = []
    for ln in lines:
        while ln and fixed and ln[0] in NO_START:
            fixed[-1] += ln[0]; ln = ln[1:]
        fixed.append(ln)
    return [l for l in fixed if l] or [text]


def rebuild_with_breaks(paragraph, lines, size_pt=None):
    """段落を、指定行に <a:br> で区切って作り直す（書式は元の1つ目のrunを継承）。"""
    p = paragraph._p
    orig = paragraph.runs[0]._r.find(qn("a:rPr")) if paragraph.runs else None
    for el in list(p):
        if el.tag in (qn("a:r"), qn("a:br")):
            p.remove(el)
    endpr = p.find(qn("a:endParaRPr"))

    def mk_run(txt):
        r = p.makeelement(qn("a:r"), {})
        rPr = copy.deepcopy(orig) if orig is not None else r.makeelement(qn("a:rPr"), {})
        if size_pt:
            rPr.set("sz", str(int(size_pt * 100)))
        r.append(rPr)
        t = r.makeelement(qn("a:t"), {}); t.text = txt; r.append(t)
        return r

    nodes = []
    for i, ln in enumerate(lines):
        if i:
            nodes.append(p.makeelement(qn("a:br"), {}))
        nodes.append(mk_run(ln))
    for nd in nodes:
        (endpr.addprevious(nd) if endpr is not None else p.append(nd))


def add_text_shadow(shape, alpha=72, blur=4, dist=3, direction=2700000, color="000000"):
    """図形（＝テキスト）に外側の影を付けて写真上でも視認性を上げる。"""
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"),
                         {"blurRad": str(int(blur * 12700)), "dist": str(int(dist * 12700)),
                          "dir": str(direction), "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": color})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))}))
    sh.append(clr); eff.append(sh); spPr.append(eff)


## ===== 高級バッジ／統一キャプション（再現性のための部品） =====
GOLD = "BEAF87"
GOLD_STOPS = [(0, "6E4F22", 100), (22, "F3E6BC", 100), (48, "B9975B", 100),
              (74, "F7EFD2", 100), (100, "7A5A2C", 100)]   # メタリック箔
GLASS_STOPS = [(0, "2A2F38", 82), (55, "161A20", 80), (100, "0B0D11", 84)]  # スモークガラス


def _grad_el(parent, stops, ang):
    gf = parent.makeelement(qn("a:gradFill"), {}); lst = gf.makeelement(qn("a:gsLst"), {})
    for pos, col, al in stops:
        gs = lst.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        c = gs.makeelement(qn("a:srgbClr"), {"val": col})
        if al < 100:
            c.append(c.makeelement(qn("a:alpha"), {"val": str(int(al * 1000))}))
        gs.append(c); lst.append(gs)
    gf.append(lst); gf.append(gf.makeelement(qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
    return gf


def _shape_grad(shape, stops, ang):
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    geom = spPr.find(qn("a:prstGeom")); gf = _grad_el(spPr, stops, ang)
    (geom.addnext(gf) if geom is not None else spPr.append(gf))


def _bevel(shape, w=5, h=3):
    spPr = shape._element.spPr
    s = spPr.makeelement(qn("a:sp3d"), {})
    s.append(s.makeelement(qn("a:bevelT"), {"w": str(int(w * 12700)), "h": str(int(h * 12700)), "prst": "circle"}))
    spPr.append(s)


def _shadow(shape, blur=8, dist=4, direction=5400000, a=52, color="000000"):
    spPr = shape._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {"blurRad": str(int(blur * 12700)), "dist": str(int(dist * 12700)),
                                             "dir": str(direction), "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": color}); clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(a * 1000))}))
    sh.append(clr); eff.append(sh); spPr.append(eff)


def _alpha(fore, pct):
    e = fore._xFill.find(qn("a:srgbClr"))
    if e is not None:
        e.append(e.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))


def _rrect(slide, x, y, w, h, rad):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x * CM), int(y * CM), int(w * CM), int(h * CM))
    try:
        s.adjustments[0] = rad
    except Exception:
        pass
    s.line.fill.background(); s.shadow.inherit = False
    return s


def _bar(slide, x, y, w, h, col):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x * CM), int(y * CM), int(w * CM), int(h * CM))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(col); s.line.fill.background(); s.shadow.inherit = False
    return s


def _txt(slide, x, y, w, h, items):
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tb = slide.shapes.add_textbox(int(x * CM), int(y * CM), int(w * CM), int(h * CM))
    tf = tb.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, sz, col, bold, spc, grad) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(0); p.space_after = Pt(0)
        r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = bold; r.font.name = "HGS明朝E"
        if col:
            r.font.color.rgb = RGBColor.from_string(col)
        rPr = r._r.get_or_add_rPr(); rPr.set("spc", str(int(spc)))
        for tag in ("a:latin", "a:ea", "a:cs"):
            e = rPr.find(qn(tag))
            if e is None:
                e = rPr.makeelement(qn(tag), {}); rPr.append(e)
            e.set("typeface", "HGS明朝E")
        if grad:
            rPr2 = r._r.get_or_add_rPr()
            for tg in ("a:noFill", "a:solidFill", "a:gradFill"):
                for e in rPr2.findall(qn(tg)):
                    rPr2.remove(e)
            gf = _grad_el(rPr2, GOLD_STOPS, 90); anc = rPr2.find(qn("a:latin"))
            (anc.addprevious(gf) if anc is not None else rPr2.append(gf))
    return tb


def build_luxury_badge(slide, spec, byid):
    """物件情報バッジを毎回同一の高級意匠で生成する。スペック値だけ差し替え。
    spec: {remove_id, x,y,w,h, seal, madori, area, sub}
    戻り値: バッジ矩形(EMU) （キャプションのヒーロー配置に使う）"""
    rid = spec.get("remove_id")
    if rid is not None and byid.get(rid) is not None:
        byid[rid]._element.getparent().remove(byid[rid]._element)
    X, Y, W, H = spec.get("x", 0.7), spec.get("y", 6.05), spec.get("w", 4.7), spec.get("h", 3.6)
    cx = X + W / 2
    plate = _rrect(slide, X, Y, W, H, 0.10); _shape_grad(plate, GOLD_STOPS, 55); _bevel(plate); _shadow(plate)
    plate.name = "LUXBADGE_PLATE"
    glass = _rrect(slide, X + 0.11, Y + 0.11, W - 0.22, H - 0.22, 0.09); _shape_grad(glass, GLASS_STOPS, 90)
    o, ln, th, GC = 0.28, 0.4, 0.033, "E9D9A8"
    for (ox, oy, dx, dy) in [(X + o, Y + o, 1, 1), (X + W - o, Y + o, -1, 1),
                             (X + o, Y + H - o, 1, -1), (X + W - o, Y + H - o, -1, -1)]:
        _bar(slide, min(ox, ox + dx * ln), oy - th / 2, ln, th, GC)
        _bar(slide, ox - th / 2, min(oy, oy + dy * ln), th, ln, GC)
    if spec.get("seal"):
        pill = _rrect(slide, cx - 1.45, Y + 0.26, 2.9, 0.56, 0.5); _shape_grad(pill, GOLD_STOPS, 90)
        _bevel(pill, 3, 2); _shadow(pill, blur=5, dist=2, a=45)
        _txt(slide, cx - 1.45, Y + 0.28, 2.9, 0.52, [(spec["seal"], 10.5, "241A06", True, 120, False)])
    t1 = _txt(slide, X, Y + 1.0, W, 0.9, [(spec.get("madori", ""), 20, None, True, 300, True)]); _shadow(t1, blur=3, dist=2, a=60)
    _txt(slide, X, Y + 1.92, W, 0.5, [(spec.get("area", ""), 13.5, "FFFFFF", False, 60, False)])
    _bar(slide, cx - 0.78, Y + 2.52, 1.56, 0.03, "D9C48A")
    if spec.get("sub"):
        _txt(slide, X, Y + 2.62, W, 0.7, [(spec["sub"], 13, "EBDBA6", True, 200, False)])
    return (int(X * CM), int(Y * CM), int(W * CM), int(H * CM))


def unify_captions(slide, spec, byid, pics, badge_box=None, hero_box=None):
    """写真キャプションを固定サイズ＋統一スタイル（スモークガラス＋金の極細罫）に。
    文字数が変わっても枠は不変。枠写真の左下へ一定オフセットで自動整列。
    ヒーロー(大)写真のタグは hero_box の左下へ置く（バッジと被らない）。"""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    FW = int(spec.get("w_cm", 2.35) * CM); FH = int(spec.get("h_cm", 0.52) * CM); OFF = int(spec.get("off_cm", 0.16) * CM)
    hero_id = spec.get("hero_id")
    for cid in spec.get("ids", []):
        c = byid.get(cid)
        if c is None:
            continue
        best = None
        for p in pics:
            if p.left - int(0.6 * CM) <= c.left <= p.left + p.width and p.top <= c.top <= p.top + p.height + int(0.6 * CM):
                if best is None or abs((p.top + p.height) - (c.top + c.height)) < abs((best.top + best.height) - (c.top + c.height)):
                    best = p
        if cid == hero_id and hero_box is not None:      # ヒーロー写真の左下（バッジ回避）
            c.left = hero_box[0] + OFF; c.top = hero_box[1] + hero_box[3] - FH - OFF
        elif best is not None:
            c.left = best.left + OFF; c.top = best.top + best.height - FH - OFF
        elif cid == hero_id and badge_box is not None:
            c.left = badge_box[0]; c.top = badge_box[1] + badge_box[3] + int(0.14 * CM)
        else:
            c.top = c.top + c.height - FH
        c.width = FW; c.height = FH
        tf = c.text_frame; tf.auto_size = MSO_AUTO_SIZE.NONE; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(0)
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for r in para.runs:
                r.font.size = Pt(spec.get("font_pt", 8)); r.font.color.rgb = RGBColor.from_string("FFFFFF"); r.font.bold = True
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string(spec.get("fill", "14161B"))
        _alpha(c.fill.fore_color, spec.get("alpha", 68))
        if spec.get("gold_line", True):
            c.line.color.rgb = RGBColor.from_string(GOLD); c.line.width = Pt(0.5)
        c.shadow.inherit = False


def add_point_panel(slide, spec, byid):
    """view_heroモードでPOINTが写真に乗って読みづらいとき、背後に半透明の暗パネルを敷き
    文字を明色にして視認性を上げる（キャッチのスクリムと同じ発想・世界観）。"""
    x, y, w, h = spec["box_cm"]
    panel = _rrect(slide, x, y, w, h, spec.get("radius", 0.05))
    panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor.from_string(spec.get("color", "12141A"))
    _alpha(panel.fill.fore_color, spec.get("alpha", 55))
    if spec.get("gold_line", True):
        panel.line.color.rgb = RGBColor.from_string(GOLD); panel.line.width = Pt(0.5)
    else:
        panel.line.fill.background()
    panel.shadow.inherit = False
    # z順：POINT文字の直下（写真の上・文字の下）へ
    parent = None; idxs = []
    for tid in spec.get("front_ids", []):
        if byid.get(tid) is not None:
            el = byid[tid]._element; parent = el.getparent(); idxs.append(list(parent).index(el))
    if parent is not None and idxs:
        e = panel._element; e.getparent().remove(e); parent.insert(min(idxs), e)
    tc = spec.get("text_color")
    if tc:
        for tid in spec.get("text_ids", []):
            sh = byid.get(tid)
            if sh is not None:
                for para in sh.text_frame.paragraphs:
                    for r in para.runs:
                        r.font.color.rgb = RGBColor.from_string(tc)


def fit_region(region_cm, iw, ih, align="center"):
    """配置可能エリア内に、画像比率を保った最大矩形を計算して返す(EMU)。"""
    L, T, W, H = region_cm
    Ri, Rr = iw / ih, W / H
    if Ri > Rr:
        w, h = W, W / Ri
    else:
        h, w = H, H * Ri
    if align == "topright":
        x, y = L + (W - w), T
    elif align == "topleft":
        x, y = L, T
    else:
        x, y = L + (W - w) / 2, T + (H - h) / 2
    return tuple(int(v * CM) for v in (x, y, w, h))


def all_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:
            yield from all_shapes(sh.shapes)


def img_wh(path):
    with Image.open(path) as im:
        return im.size


def replace_image(slide, shp, path, mode, box=None):
    """shp(旧図形)の位置(またはbox指定)に画像を差し込み、同じzオーダーに戻す。
    box=(L,T,W,H) EMU を渡すと、その矩形を枠にする（間取図の拡大などに使用）。"""
    if box:
        L, T, W, H = box
    else:
        L, T, W, H = shp.left, shp.top, shp.width, shp.height
    sp = shp._element
    parent = sp.getparent()
    idx = list(parent).index(sp)
    iw, ih = img_wh(path)
    Ri, Rf = iw / ih, W / H
    if mode == "exact":       # 既に枠比率に合わせた画像を、枠いっぱいに（クロップ無し）
        pic = slide.shapes.add_picture(path, L, T, W, H)
    elif mode == "fit":
        if Ri > Rf:
            w = W; h = int(W / Ri)
        else:
            h = H; w = int(H * Ri)
        left = L + (W - w) // 2; top = T + (H - h) // 2
        pic = slide.shapes.add_picture(path, left, top, w, h)
    else:  # fill = 中央クロップ
        pic = slide.shapes.add_picture(path, L, T, W, H)
        if Ri > Rf:
            c = (1 - Rf / Ri) / 2; pic.crop_left = c; pic.crop_right = c
        else:
            c = (1 - Ri / Rf) / 2; pic.crop_top = c; pic.crop_bottom = c
    newel = pic._element
    parent.remove(newel)
    parent.remove(sp)
    parent.insert(idx, newel)
    return pic


# テンプレの写真枠（見本写真が入ったまま出荷されると別物件の写真が載る）
PHOTO_FRAMES = {116: "メイン写真", 32: "写真", 33: "写真", 119: "写真",
                120: "写真", 124: "写真", 128: "写真", 139: "間取図"}


def guard_sample_photos(slide, cfg, allow_missing):
    """images.json で差し替えなかった枠には**テンプレの見本写真がそのまま残る**。
    別物件の写真を載せた図面が出来てしまうのがいちばん重い事故なので既定では止める。
    --allow-missing のときは残骸を消して穴にする（他人の写真を出すよりは穴のほうがまし）。"""
    used = {int(k) for k in cfg.get("images", {})} | {int(h) for h in cfg.get("hide_ids", [])}
    if cfg.get("fullbleed_bg"):
        used.add(116)
    left = [(sid, PHOTO_FRAMES[sid]) for sid in PHOTO_FRAMES if sid not in used]
    byid = {sh.shape_id: sh for sh in all_shapes(slide.shapes)}
    left = [(sid, lb) for sid, lb in left if byid.get(sid) is not None]
    if not left:
        return
    listing = "、".join(f"id{sid}（{lb}）" for sid, lb in left)
    if not allow_missing:
        sys.exit(f"!! 見本写真が残る枠があります: {listing}\n"
                 f"   images.json でこの物件の写真を割り当てるか、使わない枠は "
                 f'"hide_ids" に入れてください。\n'
                 f"   承知のうえで空のまま進めるなら --allow-missing。")
    for sid, lb in left:
        byid[sid]._element.getparent().remove(byid[sid]._element)
    print(f"⚠ 写真未割当のため見本写真を削除しました（枠は空）: {listing}")


def main():
    src, cfgfile, out = sys.argv[1], sys.argv[2], sys.argv[3]
    allow_missing = "--allow-missing" in sys.argv[4:]
    cfg = json.load(open(cfgfile, encoding="utf-8"))
    prs = Presentation(src)
    slide = prs.slides[0]
    byid = {sh.shape_id: sh for sh in all_shapes(slide.shapes)}

    # 1) ラベル差し替え
    for sid, cells in cfg.get("relabels", {}).items():
        tbl = byid[int(sid)].table
        for rc, text in cells.items():
            r, c = map(int, rc.split(","))
            cell = tbl.cell(r, c)
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = text
                for extra in p.runs[1:]:
                    extra._r.getparent().remove(extra._r)
            else:
                p.add_run().text = text
        print(f"relabeled table {sid}: {cells}")

    # 2) 画像差し込み
    tmpdir = os.path.join(os.path.dirname(os.path.abspath(out)), "_cmp")
    qc = []   # プリフライトQCの記録
    placed = {}   # 元id -> 差し込んだ新しい画像shape
    floorplan_box = None   # 間取図の実寸（バッジのスナップ用）

    # 1.5) 背景モード=view_hero：眺望写真を全面ブリード背景として最背面に敷く
    fb = cfg.get("fullbleed_bg")
    if fb:
        L, T, W, H = fb["region_cm"]
        fb_photo = resolve_path(fb["photo"], out, allow_missing)
        if fb_photo and fb.get("stretch"):   # 比率無視で範囲いっぱいに引き伸ばす（眺望メイン向け）
            bg = slide.shapes.add_picture(fb_photo, int(L * CM), int(T * CM), int(W * CM), int(H * CM))
            print(f"fullbleed_bg(stretch) <- {os.path.basename(fb_photo)}")
        elif fb_photo:
            cf, method, loss, dpi = smart_composite(fb_photo, W, H, fb.get("kind", "view"), tmpdir)
            bg = slide.shapes.add_picture(cf, int(L * CM), int(T * CM), int(W * CM), int(H * CM))
            print(f"fullbleed_bg <- {os.path.basename(fb_photo)} [{method} loss{loss}% dpi{dpi}]")
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        # nvGrpSpPr, grpSpPr の直後(=最背面)に差し込む
        idx = 0
        for i, ch in enumerate(spTree):
            if ch.tag.endswith("}grpSpPr"):
                idx = i + 1; break
        spTree.insert(idx, bg._element)

    # 1.6) 不要図形を消す（view_heroでは枠ヒーローを背景が兼ねる 等）
    for hid in cfg.get("hide_ids", []):
        sh = byid.get(int(hid))
        if sh is not None:
            sh._element.getparent().remove(sh._element)
            print(f"hide id={hid}")
    for sid, spec in cfg.get("images", {}).items():
        shp = byid[int(sid)]
        path = resolve_path(spec["path"], out, allow_missing)
        if path is None:
            print(f"⚠ 画像なし id={sid}: {spec['path']}")
            qc.append((sid, "欠落", 0, 0, spec["path"]))
            continue
        mode = spec.get("mode", "fill")
        border = spec.get("border")

        if spec.get("region_cm"):                     # 間取図など：エリア内で最大化
            iw, ih = img_wh(path)
            box = fit_region(spec["region_cm"], iw, ih, spec.get("align", "center"))
            pic = replace_image(slide, shp, path, "exact", box)
            if spec.get("is_floorplan"):
                floorplan_box = box
            print(f"placed id={sid} <- {os.path.basename(path)} (region最大化 {round(box[2]/CM,1)}x{round(box[3]/CM,1)}cm)")
        elif mode == "smart":                          # 自動ハイブリッド充填
            wcm, hcm = shp.width / CM, shp.height / CM
            cf, method, loss, dpi = smart_composite(path, wcm, hcm, spec.get("kind", "room"), tmpdir)
            pic = replace_image(slide, shp, cf, "exact")
            flag = " ⚠低解像(印刷粗)" if dpi < 150 else ""
            print(f"placed id={sid} <- {os.path.basename(path)} [{method} 損失{loss}% dpi{dpi}]{flag}")
            qc.append((sid, method, loss, dpi, os.path.basename(path)))
        else:
            box = tuple(int(v * CM) for v in spec["box_cm"]) if spec.get("box_cm") else None
            pic = replace_image(slide, shp, path, mode, box)
            print(f"placed id={sid} <- {os.path.basename(path)} ({mode})")

        if pic is not None:
            placed[int(sid)] = pic
        if border and pic is not None:                 # 間取図などに細枠
            pic.line.color.rgb = RGBColor.from_string(border.get("color", "C8C8C8"))
            pic.line.width = Pt(border.get("pt", 0.75))

    if qc:
        print("\n── プリフライトQC（写真）──")
        for sid, method, loss, dpi, name in qc:
            mark = "⚠" if (dpi and dpi < 150) or loss > 60 else " "
            print(f" {mark} id={sid:<4} {method:<10} 損失{loss:>3}%  dpi{dpi:>4}  {name}")

    # 2.45) 図形のリサイズ（メインキャッチを間取図に被らせない／POINT箱を広げる等）
    #        ※ jp_wrap(折返し) や point_fix(高さ) より前に実行して、新しい幅で計算させる。
    for rz in cfg.get("resize", []):
        sh = byid.get(int(rz["id"]))
        if sh is not None:
            if "x_cm" in rz: sh.left = int(rz["x_cm"] * CM)
            if "y_cm" in rz: sh.top = int(rz["y_cm"] * CM)
            if "w_cm" in rz: sh.width = int(rz["w_cm"] * CM)
            if "h_cm" in rz: sh.height = int(rz["h_cm"] * CM)
            print(f"resize id={rz['id']} w={rz.get('w_cm')} h={rz.get('h_cm')}")

    # 2.46) フォントサイズ上書き（改行が語中で切れる箱を少し縮めて意味の切れ目で折り返す等）
    #        ※ jp_wrap より前に実行して、新しいサイズで折返しさせる。
    for sf in cfg.get("set_font", []):
        sh = byid.get(int(sf["id"]))
        if sh is not None and sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(sf["size_pt"])
            print(f"set_font id={sf['id']} -> {sf['size_pt']}pt")

    # 2.5) 不要な段落を削除（例: 路線が2つの物件でアクセス3行目を消す）
    for sid, idxs in cfg.get("drop_paragraphs", {}).items():
        tf = byid[int(sid)].text_frame
        paras = list(tf.paragraphs)
        for i in sorted(idxs, reverse=True):
            if 0 <= i < len(paras):
                paras[i]._p.getparent().remove(paras[i]._p)
        print(f"dropped paragraphs {idxs} from {sid}")

    # 2.6) 写真上の文字を読みやすくするスクリム（半透明の濃色帯）を文字の背面に敷く
    def set_alpha(fore, pct):
        srgb = fore._xFill.find(qn("a:srgbClr"))
        if srgb is not None:
            srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))
    for sc in cfg.get("scrims", []):
        box = [int(v * CM) for v in sc["box_cm"]]
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if sc.get("round") else MSO_SHAPE.RECTANGLE, *box)
        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor.from_string(sc.get("color", "000000"))
        set_alpha(shp.fill.fore_color, sc.get("alpha", 40))
        shp.line.fill.background()
        shp.shadow.inherit = False
        tgt = byid[sc["before_id"]]._element         # この図形の直前(=背面)に差し込む
        parent = tgt.getparent()
        el = shp._element; parent.remove(el); parent.insert(list(parent).index(tgt), el)
        print(f"scrim behind {sc['before_id']} box={sc['box_cm']} alpha={sc.get('alpha',40)}")

    # 2.7) 写真キャプションのタグを整える（小さく・白字・半透明の濃色チップ）
    cs = cfg.get("caption_style")
    if cs:
        size = Pt(cs.get("size_pt", 7))
        color = RGBColor.from_string(cs.get("color", "FFFFFF"))
        fillhex = cs.get("fill", "2B2B2B")
        alpha = cs.get("fill_alpha")  # 0-100 (%) 不透明度。省略で不透明
        for cid in cs.get("ids", []):
            sh = byid.get(int(cid))
            if sh is None:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    r.font.size = size
                    r.font.color.rgb = color
                    r.font.bold = cs.get("bold", False)
            # チップ塗り＝濃色（任意で半透明）
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor.from_string(fillhex)
            if alpha is not None:
                srgb = sh.fill.fore_color._xFill.find(qn("a:srgbClr"))
                if srgb is not None:
                    a = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
                    srgb.append(a)
            # 枠線は消す
            try:
                sh.line.fill.background()
            except Exception:
                pass
            # 語中で折り返さない＝1行固定（チップは文字幅に自動フィット）
            if cs.get("no_wrap", True):
                sh.text_frame.word_wrap = False
                sh.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        print(f"caption_style: ids={cs.get('ids')} size={cs.get('size_pt',7)}pt color={cs.get('color','FFFFFF')} fill={fillhex} alpha={alpha}")

    # 2.8) POINTを「・」箇条書きにする
    for sid in cfg.get("bullets", {}).get("ids", []):
        tf = byid[int(sid)].text_frame
        for para in tf.paragraphs:
            if para.text.strip() and not para.text.lstrip().startswith("・"):
                if para.runs:
                    para.runs[0].text = "・" + para.runs[0].text
        print(f"bullets: id={sid}")

    # 2.9) 本文を日本語ワードラップ（語中で切らない＋禁則）
    jw = cfg.get("jp_wrap")
    if jw:
        for spec in jw.get("items", []):
            sh = byid[int(spec["id"])]
            size_pt = spec.get("size_pt")
            max_cm = sh.width / CM * spec.get("inner_ratio", 0.94)
            sa = spec.get("space_after")
            for para in sh.text_frame.paragraphs:
                txt = para.text
                if not txt.strip():
                    continue
                pt = size_pt or (para.runs[0].font.size.pt if para.runs and para.runs[0].font.size else 9)
                lines = jp_wrap(txt, max_cm, pt)
                rebuild_with_breaks(para, lines, size_pt)
                if sa is not None:
                    para.space_after = Pt(sa); para.space_before = Pt(0)
            print(f"jp_wrap: id={spec['id']} size={size_pt}")

    # 3) POINTが写真にかぶらないよう固定
    pf = cfg.get("point_fix")
    if pf:
        sh = byid[pf["id"]]
        tf = sh.text_frame
        # 末尾の空段落を除去
        paras = tf.paragraphs
        for p in list(paras)[::-1]:
            if p.text.strip() == "":
                p._p.getparent().remove(p._p)
            else:
                break
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        bottom = int(pf["bottom_cm"] * CM)
        sh.height = max(int(0.5 * CM), bottom - sh.top)
        print(f"point_fix id={pf['id']}: bottom={pf['bottom_cm']}cm")

    # 4) 写真上の文字に影を付けて視認性を上げる
    for spec in cfg.get("text_shadow", []):
        for cid in spec.get("ids", []):
            sh = byid.get(int(cid))
            if sh is not None:
                add_text_shadow(sh, alpha=spec.get("alpha", 72), blur=spec.get("blur", 4),
                                dist=spec.get("dist", 3), color=spec.get("color", "000000"))
        print(f"text_shadow: ids={spec.get('ids')}")

    # 4.4) バッジ等を間取図の実寸にスナップ（間取図サイズが可変でも隙間なく整列）
    bs = cfg.get("badge_snap")
    if bs and floorplan_box is not None:
        fx, fy, fw, fh = floorplan_box
        sh = byid.get(int(bs["id"]))
        if sh is not None:
            gap = int(bs.get("gap_cm", 0.15) * CM)
            align = bs.get("align", "under_center")
            if align == "under_center":
                sh.top = fy + fh + gap
                sh.left = fx + (fw - sh.width) // 2
            elif align == "under_left":
                sh.top = fy + fh + gap; sh.left = fx
            elif align == "left_top":     # 間取図の左隣・上
                sh.left = fx - sh.width - gap; sh.top = fy
            print(f"badge_snap id={bs['id']} -> under floorplan ({align})")

    # 4.5) 図形の移動（間取図の領域から重なる図形を退避 等）
    for mv in cfg.get("move", []):
        sh = byid.get(int(mv["id"]))
        if sh is not None:
            L, T = mv["to_cm"]
            sh.left = int(L * CM); sh.top = int(T * CM)
            print(f"move id={mv['id']} -> {mv['to_cm']}")

    # 4.56) 指定図形の背面に白板（物件概要の後ろ 等）／POINTの視認性用の白半透明板にも使う
    #   box_cm=[x,y,w,h] で明示座標（上帯下〜下帯上まで伸ばす等）。無ければ behind_id の図形範囲＋pad。
    #   to_back=true で最背面へ（右カラム全体の後ろに敷く。左カラムと重ならない前提）。
    for bp in cfg.get("bg_panels", []):
        tgt = byid.get(int(bp["behind_id"])) if bp.get("behind_id") is not None else None
        if bp.get("box_cm"):
            x, y, w, h = [int(v * CM) for v in bp["box_cm"]]
        elif tgt is not None:
            pad = int(bp.get("pad_cm", 0.0) * CM)
            x, y, w, h = tgt.left - pad, tgt.top - pad, tgt.width + 2 * pad, tgt.height + 2 * pad
        else:
            continue
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if bp.get("round") else MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor.from_string(bp.get("color", "FFFFFF"))
        if bp.get("alpha") is not None:
            _alpha(shp.fill.fore_color, bp["alpha"])
        if bp.get("line_color"):
            shp.line.color.rgb = RGBColor.from_string(bp["line_color"]); shp.line.width = Pt(bp.get("line_pt", 0.75))
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        el = shp._element; parent = el.getparent()
        parent.remove(el)
        if bp.get("to_back"):                    # 最背面（grpSpPr直後）へ
            idx = 0
            for i, ch in enumerate(parent):
                if ch.tag.endswith("}grpSpPr"):
                    idx = i + 1; break
            parent.insert(idx, el)
        elif tgt is not None:                    # 対象図形の直前（＝すぐ背面）へ
            parent.insert(list(parent).index(tgt._element), el)
        print(f"bg_panel behind={bp.get('behind_id')} box={bp.get('box_cm')} to_back={bp.get('to_back')}")

    # 5) 指定図形を最前面へ（間取図の前面には何も置かない）
    for cid in cfg.get("bring_front", []):
        sh = placed.get(int(cid)) or byid.get(int(cid))
        if sh is not None:
            el = sh._element
            parent = el.getparent()
            parent.remove(el); parent.append(el)
            print(f"bring_front: id={cid}")

    # 6) 高級バッジ（毎回同一意匠・値だけ差し替え）
    badge_box = None
    if cfg.get("luxury_badge"):
        badge_box = build_luxury_badge(slide, cfg["luxury_badge"], byid)
        print("luxury_badge:", cfg["luxury_badge"].get("madori"), cfg["luxury_badge"].get("seal"), cfg["luxury_badge"].get("sub"))

    # 6.5) POINTパネル（view_heroで写真の上のPOINTを読みやすく）
    if cfg.get("point_panel"):
        add_point_panel(slide, cfg["point_panel"], byid)
        print("point_panel: applied")

    # 7) キャプション統一（固定サイズ＋金の極細罫）
    if cfg.get("caption_uniform"):
        curr = list(all_shapes(slide.shapes))
        cbyid = {s.shape_id: s for s in curr}
        pics = [s for s in curr if s.shape_type == 13 and 3 * CM < s.width < 8 * CM and s.height < 6 * CM]
        hp = placed.get(int(cfg["caption_uniform"].get("hero_photo_id", -1)))
        hero_box = (hp.left, hp.top, hp.width, hp.height) if hp is not None else None
        unify_captions(slide, cfg["caption_uniform"], cbyid, pics, badge_box, hero_box)
        print("caption_uniform:", cfg["caption_uniform"].get("ids"))

    guard_sample_photos(slide, cfg, allow_missing)

    prs.save(out)
    print("saved:", out)


if __name__ == "__main__":
    main()
